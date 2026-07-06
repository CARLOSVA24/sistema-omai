import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Paleta de Colores
    NAVY = RGBColor(11, 25, 44)       # #0B192C (Azul Marino Oscuro)
    STEEL = RGBColor(30, 62, 98)      # #1E3E62 (Azul Acero)
    GOLD = RGBColor(224, 169, 109)    # #E0A96D (Dorado Militar)
    OFFWHITE = RGBColor(245, 247, 248) # #F5F7F8 (Fondo Claro)
    WHITE = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(30, 41, 59)   # #1E293B (Slate 800)
    MUTED_TEXT = RGBColor(100, 116, 139) # #64748B (Slate 500)
    
    # ----------------------------------------------------
    # FUNCIONES AUXILIARES PARA DISEÑO DE DIAPOSITIVAS Y TRANSICIONES
    # ----------------------------------------------------
    
    def apply_transition(slide, t_type="fade"):
        """Inyecta XML de transición para dar movimiento a las diapositivas."""
        if t_type == "fade":
            xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>'
        elif t_type == "push_left":
            xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>'
        elif t_type == "push_down":
            xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="dn"/></p:transition>'
        elif t_type == "wipe_right":
            xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe dir="r"/></p:transition>'
        elif t_type == "zoom":
            xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:zoom/></p:transition>'
        else:
            return
        
        try:
            xml_fragment = parse_xml(xml)
            slide.element.insert(-1, xml_fragment)
        except Exception as e:
            print(f"Error inyectando transición: {e}")

    def add_title_slide(prs, title, subtitle, author, institution):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Fondo Azul Marino
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY
        
        # Triángulo Decorativo Superior Derecho (Azul Acero)
        shape1 = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.333), Inches(0), Inches(4.0), Inches(4.0)
        )
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = STEEL
        shape1.line.color.rgb = STEEL
        shape1.rotation = 90
        
        # Triángulo Decorativo Superior Derecho (Dorado)
        shape2 = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE, Inches(10.333), Inches(0), Inches(3.0), Inches(3.0)
        )
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = GOLD
        shape2.line.color.rgb = GOLD
        shape2.rotation = 90
        
        # Imagen del escudo de la Armada (si existe)
        logo_path = "ESCUDOARMADA.jpg"
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(1.5), Inches(1.2), width=Inches(1.5), height=Inches(1.5))
            
        # Cuadro de Título
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.9), Inches(10.3), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.name = "Segoe UI"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = GOLD
        
        # Cuadro de Subtítulo
        sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.8))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = WHITE
        
        # Línea de División Dorada
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.2), Inches(8.0), Inches(0.04)
        )
        div.fill.solid()
        div.fill.fore_color.rgb = GOLD
        div.line.color.rgb = GOLD
        
        # Cuadro de Autoría e Institución
        auth_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.4), Inches(10.3), Inches(1.2))
        tf_auth = auth_box.text_frame
        tf_auth.word_wrap = True
        
        p_inst = tf_auth.paragraphs[0]
        p_inst.text = institution
        p_inst.font.name = "Segoe UI"
        p_inst.font.size = Pt(14)
        p_inst.font.bold = True
        p_inst.font.color.rgb = GOLD
        p_inst.space_after = Pt(4)
        
        p_auth = tf_auth.add_paragraph()
        p_auth.text = author
        p_auth.font.name = "Segoe UI"
        p_auth.font.size = Pt(12)
        p_auth.font.color.rgb = WHITE
        
        # Transición de Portada (Fade suave al iniciar)
        apply_transition(slide, "fade")
        
    def add_section_slide(prs, title, subtitle):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Fondo Azul Marino
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY
        
        # Banda Lateral Dorada
        decor = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
        )
        decor.fill.solid()
        decor.fill.fore_color.rgb = GOLD
        decor.line.color.rgb = GOLD
        
        # Banda Lateral Azul Acero
        decor2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0), Inches(0.15), Inches(7.5)
        )
        decor2.fill.solid()
        decor2.fill.fore_color.rgb = STEEL
        decor2.line.color.rgb = STEEL
        
        # Título de Sección
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.name = "Segoe UI"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = GOLD
        
        # Subtítulo de Sección
        sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.5))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(22)
        p_sub.font.color.rgb = WHITE
        
        # Transición de Sección: Desplazamiento Izquierda (Push Left)
        apply_transition(slide, "push_left")
        
    def add_content_slide(prs, title_text, bullets, image_path, footer_text, transition="fade"):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Fondo Claro
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = OFFWHITE
        
        # Encabezado (Banda Superior Azul Marino)
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = NAVY
        header_shape.line.color.rgb = NAVY
        
        # Título de la Diapositiva
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = "Segoe UI"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.LEFT
        
        # Línea de División Dorada
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = GOLD
        divider.line.color.rgb = GOLD
        
        # Columna Izquierda: Explicación y Viñetas
        text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(5.6), Inches(5.3))
        tf_content = text_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_left = Inches(0)
        tf_content.margin_right = Inches(0)
        
        for idx, (head, desc) in enumerate(bullets):
            p = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
            p.space_after = Pt(4)
            p.space_before = Pt(8)
            p.line_spacing = 1.15
            
            # Encabezado de la viñeta
            run_head = p.add_run()
            run_head.text = "• " + head + ": "
            run_head.font.name = "Segoe UI"
            run_head.font.size = Pt(14)
            run_head.font.bold = True
            run_head.font.color.rgb = STEEL
            
            # Descripción de la viñeta
            run_desc = p.add_run()
            run_desc.text = desc
            run_desc.font.name = "Segoe UI"
            run_desc.font.size = Pt(12)
            run_desc.font.color.rgb = DARK_TEXT
            
        # Columna Derecha: Captura de Pantalla
        if os.path.exists(image_path):
            # Sombra gris detrás de la imagen (Offset de 0.05")
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(6.55), Inches(1.45), Inches(6.2), Inches(4.8)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(220, 224, 230)
            shadow.line.color.rgb = RGBColor(220, 224, 230)
            
            # Imagen
            slide.shapes.add_picture(
                image_path, Inches(6.5), Inches(1.4), width=Inches(6.2), height=Inches(4.8)
            )
            
            # Borde Decorativo sobre la imagen
            border = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.4), Inches(6.2), Inches(4.8)
            )
            border.fill.background()
            border.line.color.rgb = STEEL
            border.line.width = Pt(1.5)
            
        # Pie de Página
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.133), Inches(0.4))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = footer_text
        p_foot.font.name = "Segoe UI"
        p_foot.font.size = Pt(10)
        p_foot.font.italic = True
        p_foot.font.color.rgb = MUTED_TEXT
        
        # Aplicar la transición indicada (por defecto Fade suave)
        apply_transition(slide, transition)
        
    # ----------------------------------------------------
    # GENERACIÓN DE DIAPOSITIVAS
    # ----------------------------------------------------
    
    # 1. Portada
    print("Creando Portada...")
    add_title_slide(
        prs,
        title="SISTEMA OMAI",
        subtitle="Plataforma de Mando, Control y Gestión Operacional para Patrullajes Terrestres",
        author="Autor: CPCB-SU Carlos Vallejo Ortega | Versión 2.0.0 (2026)",
        institution="ARMADA DEL ECUADOR - GRUPO DE TAREA 100.51"
    )
    
    # 2. Introducción y Arquitectura
    print("Creando Diapositiva 2: Introducción...")
    add_content_slide(
        prs,
        title_text="⚙️ Arquitectura y Tecnología del Sistema",
        bullets=[
            ("Servidor Local Robusto", "Empaquetado en un ejecutable local (SISTEMA_OMAI.exe) que expone la API en el puerto 3000, facilitando su ejecución offline."),
            ("Base de Datos SQLite", "Utiliza database.sqlite para un almacenamiento rápido, local y confiable, idóneo para despliegues aislados o móviles."),
            ("Conexión en Tiempo Real", "Implementa WebSockets (Socket.io) para sincronizar inmediatamente los cambios y logs entre múltiples dispositivos conectados."),
            ("Exportaciones Directas", "Integra librerías jsPDF y SheetJS para generar al instante hojas de cálculo Excel e informes oficiales listos para impresión.")
        ],
        image_path="screenshots/01_login.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Tecnología y Despliegue Local",
        transition="fade"
    )
    
    # ----------------------------------------------------
    # SECCIÓN 1: PERSONAL
    # ----------------------------------------------------
    print("Creando Sección 1: Personal...")
    add_section_slide(
        prs,
        title="MÓDULO DE PERSONAL",
        subtitle="Administración de Fuerza, Guardias, Distribución Táctica y Estadísticas Dinámicas"
    )
    
    # 4. Registro de Personal
    add_content_slide(
        prs,
        title_text="👥 Registro de Personal y Carga de Dotación",
        bullets=[
            ("Ficha de Personal Detallada", "Registra grado militar, especialidad, nombres completos, cédula, número de contacto, reparto de origen y rotación."),
            ("Carga Masiva desde Excel", "Importa listados completos con detección inteligente de columnas y guardias, reduciendo el trabajo manual."),
            ("Organización por Grupos", "Clasifica el personal por condición operacional (Operativo/Franco) y les asigna su respectivo grupo de destino (GT ECHO, CODESC, etc.)."),
            ("Búsqueda en Tiempo Real", "Buscador integrado por nombre o cédula que filtra la base de datos de inmediato en la tabla interactiva.")
        ],
        image_path="screenshots/02_registro_personal.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Personal - Registro",
        transition="fade"
    )
    
    # 5. Distribución de Personal (Guardias)
    add_content_slide(
        prs,
        title_text="👥 Planificación de Guardias y Turnos",
        bullets=[
            ("División por Guardias", "Separa rápidamente a la dotación activa en dos grupos principales: Guardia de Babor y Guardia de Estribor."),
            ("Algoritmo de Turnos", "Distribuye automáticamente al personal en turnos rotativos tradicionales de guardia (T1, T2 y T3) de manera equilibrada."),
            ("Exportaciones Tácticas", "Botones dedicados para exportar el cuadro general de guardias a PDF estructurado u hojas de cálculo Excel."),
            ("Régimen 2x2 CODESC", "Aplica un régimen cíclico automatizado de 2 días de trabajo por 2 días de franco, recalculando la disponibilidad según fecha.")
        ],
        image_path="screenshots/03_distribucion_personal.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Personal - Guardias y Turnos",
        transition="fade"
    )
    
    # 6. Distribución por Puestos
    add_content_slide(
        prs,
        title_text="👥 Distribución Horaria por Puestos Físicos",
        bullets=[
            ("Cuadrante de Puestos", "Asigna al personal en turnos horarios exactos en los puestos físicos configurados del recinto militar."),
            ("Gestión de Apoyos y Cuotas", "Asigna tareas de apoyo complementario que requieren cuotas específicas de dotación sin entorpecer los turnos principales."),
            ("Control Visual Rápido", "Presenta un cuadrante limpio y ordenado que facilita al Oficial de Guardia el control de las garitas y puntos de control."),
            ("Personalización Táctica", "Permite reasignar puestos o realizar cambios sobre la marcha para responder a novedades del servicio.")
        ],
        image_path="screenshots/04_distribucion_puestos.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Personal - Distribución por Puestos",
        transition="fade"
    )
    
    # 7. Registro de Puesto de Mando
    add_content_slide(
        prs,
        title_text="👥 Registro de Puesto de Mando",
        bullets=[
            ("Oficiales de Guardia", "Permite registrar al personal al mando de la guardia del día, incluyendo Oficial de Guardia (ODG) y auxiliares."),
            ("Integración con Documentos", "Los oficiales registrados en el Puesto de Mando firman de forma automática los partes diarios y aprueban patrullajes."),
            ("Control Centralizado", "Establece claramente los responsables del monitoreo de incidentes y las comunicaciones durante la guardia.")
        ],
        image_path="screenshots/05_puesto_mando.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Personal - Puesto de Mando",
        transition="fade"
    )
    
    # 8. Designación Otras Funciones
    add_content_slide(
        prs,
        title_text="👥 Otras Funciones y Roles Especiales",
        bullets=[
            ("Asignaciones Asimétricas", "Registra al personal destinado a tareas no convencionales (Rancheros, Conductores de Ambulancia, Enfermeros, etc.)."),
            ("Prevención de Conflictos", "Separa a estos elementos del cuadrante regular de turnos automáticos, evitando que un marinero sea asignado a dos roles."),
            ("Visualización de Excepciones", "Muestra de forma segregada a la dotación que realiza funciones de apoyo permanente o administrativo.")
        ],
        image_path="screenshots/06_otras_funciones.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Personal - Roles Especiales",
        transition="fade"
    )
    
    # 9. Estadísticas de Personal
    add_content_slide(
        prs,
        title_text="👥 Estadísticas Dinámicas y Capacidad Operativa",
        bullets=[
            ("Indicadores KPI de Fuerza", "Muestra el total registrado, la cifra de personal operativo real y la cantidad de elementos de franco o en guardia."),
            ("Gráficos de Repartos y Grados", "Presenta gráficos apilados dinámicos con la distribución de rangos militares asignados a cada reparto de la dotación."),
            ("Chart.js Integrado", "Gráficos vectoriales e interactivos que responden al pasar el cursor, ideales para informes de personal."),
            ("Justificación de Recursos", "Genera de forma visual reportes instantáneos para justificar estados de fuerza ante los comandos superiores.")
        ],
        image_path="screenshots/07_estadisticas_personal.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Personal - Estadísticas",
        transition="fade"
    )
    
    # ----------------------------------------------------
    # SECCIÓN 2: OPERACIONES
    # ----------------------------------------------------
    print("Creando Sección 2: Operaciones...")
    add_section_slide(
        prs,
        title="MÓDULO DE OPERACIONES",
        subtitle="Planificación Diaria, Gestión de Órdenes de Patrulla, Plantillas y Partes Oficiales"
    )
    
    # 11. Planificación Diaria
    add_content_slide(
        prs,
        title_text="🛡️ Planificación Diaria de Operaciones",
        bullets=[
            ("Cronograma de Eventos", "Permite planificar las actividades operacionales críticas del día con horarios, sectores y personal asignado."),
            ("Agenda del Puesto de Mando", "Mantiene al Oficial de Guardia alineado con las directrices operativas dispuestas por el mando táctico."),
            ("Trazabilidad Operacional", "Registra hitos y actividades especiales (inspecciones, visitas, adiestramientos) dentro de una única agenda integrada.")
        ],
        image_path="screenshots/08_planificacion_operaciones.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Operaciones - Planificación",
        transition="fade"
    )
    
    # 12. Órdenes de Patrulla GT ECHO
    add_content_slide(
        prs,
        title_text="🛡️ Órdenes de Patrulla Estandarizadas - GT ECHO",
        bullets=[
            ("Doctrina Militar", "Formulario estructurado bajo el formato formal de cinco párrafos de la Armada: Situación, Misión, Ejecución, Logística y Enlace."),
            ("Campos Autocompletados", "Pre-llena datos automáticos como Grupo de Fecha y Hora (DTG) del sistema y correlativos de órdenes anteriores."),
            ("Vínculo de Personal", "Asocia la patrulla con el Oficial de Guardia (ODG) registrado y las dotaciones en turno de la Guardia."),
            ("Diseño de Impresión Formal", "Genera el documento final en un PDF formal, incluyendo marcas de clasificación de seguridad y bloque de firmas.")
        ],
        image_path="screenshots/09_ordpat_gt_echo.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Operaciones - Orden GT ECHO",
        transition="fade"
    )
    
    # 13. Órdenes de Patrulla GT 100.51
    add_content_slide(
        prs,
        title_text="🛡️ Órdenes de Patrulla - GT 100.51",
        bullets=[
            ("Personalización del Reparto", "Formato adaptado con encabezado oficial y logotipos específicos para el Grupo de Tarea 100.51 de Seguridad Marítima."),
            ("Precisión en Coordinaciones", "Secciones detalladas para coordinar frecuencias de radio, sectores de patrullaje, reglas de uso de la fuerza y enlaces."),
            ("Seguridad del Documento", "Clasificación de seguridad automática ('SECRETO' / 'CONFIDENCIAL') según las directrices operacionales."),
            ("Descarga Directa", "Permite la impresión física del documento doctrinario oficial de operaciones en tamaño A4 directamente desde el navegador.")
        ],
        image_path="screenshots/10_ordpat_gt_100_51.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Operaciones - Orden GT 100.51",
        transition="fade"
    )
    
    # 14. Registro de Gestión de Órdenes y Plantillas
    add_content_slide(
        prs,
        title_text="🛡️ Gestión de Plantillas Operacionales",
        bullets=[
            ("Biblioteca de Órdenes", "Consolida un listado de todas las órdenes pre-guardadas o generadas previamente para consulta o reedición."),
            ("Clonado Eficiente de Órdenes", "Función 'Pre-llenar con la última' que recupera toda la estructura doctrinaria en 1 segundo, requiriendo solo cambiar la fecha."),
            ("Ahorro de Tiempo de Oficina", "Reduce drásticamente el tiempo de elaboración administrativa de órdenes tácticas repetitivas de 15 minutos a 30 segundos."),
            ("Control de Estados de Firma", "Muestra de forma jerárquica cuáles órdenes han sido redactadas, cuáles están aprobadas y cuáles están archivadas.")
        ],
        image_path="screenshots/11_registro_gestion_ordenes.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Operaciones - Plantillas",
        transition="fade"
    )
    
    # 15. Partes al Instante
    add_content_slide(
        prs,
        title_text="🛡️ Partes al Instante y Correspondencia",
        bullets=[
            ("Cinta de Formato Estilo Word", "Editor de texto enriquecido integrado que permite aplicar negrita, cursiva, subrayado, tamaños y colores a los informes."),
            ("Simulación de Hoja Física A4", "Diseño en pantalla idéntico a un papel membretado impreso institucional con la clasificación en rojo 'SECRETO'."),
            ("Firma del Oficial de Guardia", "Vínculo automático del firmante de guardia registrado, permitiendo redactar partes tácticos urgentes y prolijos."),
            ("Exportación Directa a PDF", "Genera el PDF con la tipografía y el espaciado exacto estandarizado por las normas de correspondencia de la Armada.")
        ],
        image_path="screenshots/12_partes_al_instante.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Operaciones - Partes al Instante",
        transition="fade"
    )
    
    # 16. Órdenes de Operación
    add_content_slide(
        prs,
        title_text="🛡️ Repositorio de Órdenes de Operación",
        bullets=[
            ("Centralización de Directrices", "Permite cargar las órdenes generales de operación (OROP) y directivas emitidas por escalones superiores."),
            ("Consulta Rápida de Campo", "El personal del Puesto de Mando puede consultar en segundos las directrices operacionales vigentes para la misión."),
            ("Acceso Seguro por Roles", "Restricción de lectura para asegurar que solo oficiales autorizados visualicen directivas de alto nivel.")
        ],
        image_path="screenshots/13_ordenes_operacion.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Operaciones - Órdenes de Operación",
        transition="fade"
    )
    
    # ----------------------------------------------------
    # SECCIÓN 3: LOGÍSTICA E INTELIGENCIA
    # ----------------------------------------------------
    print("Creando Sección 3: Logística e Inteligencia...")
    add_section_slide(
        prs,
        title="LOGÍSTICA E INTELIGENCIA",
        subtitle="Control de Medios de Movilidad Terrestre y Análisis Geoespacial Avanzado del Delito"
    )
    
    # 18. Registro de Vehículos (Logística)
    add_content_slide(
        prs,
        title_text="📦 Logística: Control de Flota y Conductores",
        bullets=[
            ("Inventario de Vehículos", "Administra las unidades asignadas (camiones tácticos, camionetas de patrulla, cuadrones) ingresando placa y tipo."),
            ("Asignación de Conductores", "Vincula choferes militares específicos a cada vehículo, estableciendo la responsabilidad directa del medio."),
            ("Disponibilidad y Combustible", "Controla el estado operativo de los vehículos (Operativo / Fuera de Servicio) y registra el nivel de combustible."),
            ("Planificación Logística", "Garantiza que el Puesto de Mando asigne patrullas únicamente a medios terrestres listos y abastecidos.")
        ],
        image_path="screenshots/14_registro_vehiculos.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Logística - Vehículos",
        transition="fade"
    )
    
    # 19. Inteligencia - Mapa y Herramientas de Dibujo
    add_content_slide(
        prs,
        title_text="🧠 Inteligencia: Análisis Geoespacial",
        bullets=[
            ("Mapa Leaflet Interactivo", "Visualizador geográfico centrado en el área de operaciones con marcadores de delitos y cámaras de seguridad."),
            ("Herramientas de Dibujo Geográfico", "Permite dibujar polígonos (sectores de patrullaje), líneas (rutas de patrulla) y círculos (áreas de búsqueda)."),
            ("Carga de Capas KMZ/KML", "Permite importar archivos KMZ de Google Earth para superponer límites distritales o áreas críticas oficiales en el mapa."),
            ("Captura de Coordenadas", "Asigna latitud y longitud exactas a cada delito con un solo clic sobre el mapa para poblar el parte del incidente.")
        ],
        image_path="screenshots/15_inteligencia_mapa.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Inteligencia - Visualizador",
        transition="fade"
    )
    
    # 20. Inteligencia - Tabla de Novedades e Incidentes
    add_content_slide(
        prs,
        title_text="🧠 Inteligencia: Consolidado e Historial delictivo",
        bullets=[
            ("Base de Datos delictiva", "Registro de delitos críticos: Sicariatos, Robos, Porte de Armas, Atentados, Narcotráfico y Cámaras de Seguridad."),
            ("Filtros de Búsqueda Dinámicos", "Permite filtrar de inmediato toda la tabla por Tipo de Delito o Distrito (SUR, ESTEROS, 9 DE OCTUBRE, MODELO)."),
            ("Exportaciones a Excel/PDF", "Permite descargar el historial completo de incidentes directamente para análisis o auditorías externas."),
            ("Edición y Eliminación", "Control completo sobre cada registro de incidente para corregir descripciones o coordenadas erróneas.")
        ],
        image_path="screenshots/16_inteligencia_tabla.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Inteligencia - Tabla de Incidentes",
        transition="fade"
    )
    
    # 21. Inteligencia - Infografía Estratégica
    add_content_slide(
        prs,
        title_text="🧠 Inteligencia: Infografía Estratégica",
        bullets=[
            ("KPIs Ejecutivos de Seguridad", "Indicadores de alto nivel: Total de incidentes, delito predominante, provincia crítica y horario de mayor riesgo."),
            ("Cuatro Gráficos Analíticos", "Desglose visual interactivo por provincia (barras), por ciudad (barras), tipología delictiva (pie) y tendencia temporal."),
            ("Semáforos de Alerta Territorial", "Clasifica automáticamente las zonas en alertas Roja, Amarilla o Verde según su nivel de delincuencia."),
            ("Exportación de Infografía", "Botón rápido para compilar el informe infográfico dinámico en una presentación PDF lista para el Comandante.")
        ],
        image_path="screenshots/17_inteligencia_infografia.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Módulo de Inteligencia - Infografía",
        transition="fade"
    )
    
    # ----------------------------------------------------
    # SECCIÓN 4: HISTÓRICOS Y GESTIÓN GENERAL
    # ----------------------------------------------------
    print("Creando Sección 4: Históricos y Administración...")
    add_section_slide(
        prs,
        title="HISTÓRICOS Y GESTIÓN GENERAL",
        subtitle="Trazabilidad, Resguardo de Datos, Seguridad RBAC y Bitácora de Auditoría"
    )
    
    # 23. Historial de Órdenes y Partes
    add_content_slide(
        prs,
        title_text="📜 Históricos: Registros Históricos y Trazabilidad",
        bullets=[
            ("Archivo de Órdenes de Patrulla", "Historial de todas las órdenes emitidas, con filtros para revisar directivas ejecutadas en fechas anteriores."),
            ("Partes del Instante Históricos", "Resguardo de todos los reportes de guardia firmados digitalmente para auditorías operacionales."),
            ("Reportes de Días OMAI", "Consolida las actividades tácticas diarias completadas para evaluar el cumplimiento de metas y justificar patrullajes."),
            ("Búsqueda Avanzada", "Buscador integrado para ubicar órdenes pasadas por correlativo, fecha o responsable del mando.")
        ],
        image_path="screenshots/18_historico_ordenes.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Registros Históricos",
        transition="fade"
    )
    
    # 24. Administración - Gestión de Claves y RBAC
    add_content_slide(
        prs,
        title_text="🔐 Administración: Roles y Control de Acceso (RBAC)",
        bullets=[
            ("Perfiles Diferenciados", "Controles de seguridad adaptados para ADMINISTRADOR, JEFE OMAI, PERSONAL, LOGÍSTICA, INTELIGENCIA y CMDTE GT 51."),
            ("Gestión de Claves Centralizada", "El rol de Administrador dispone de una consola para ver, editar y redefinir contraseñas de todos los perfiles."),
            ("Seguridad de la Información", "Garantiza que solo personal autorizado modifique listados de personal o acceda a la infografía criminal."),
            ("Clave Maestra de Purga", "Configuración de una contraseña especial de seguridad para proteger acciones destructivas de base de datos.")
        ],
        image_path="screenshots/20_admin_claves.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Administración - Roles y Claves",
        transition="fade"
    )
    
    # 25. Administración - Gestión de Datos
    add_content_slide(
        prs,
        title_text="🔐 Administración: Consola de Gestión de Datos",
        bullets=[
            ("Respaldos Rápidos", "Mecanismo ágil para descargar la base de datos local SQLite y resguardar la información ante fallos de hardware."),
            ("Importación de Respaldos", "Permite subir una base de datos anterior para restaurar el sistema en un equipo nuevo en segundos."),
            ("Purga Segura de Datos", "Elimina datos de personal, vehículos o delitos con protección de clave maestra para evitar borrados accidentales."),
            ("Modificaciones de Estructura", "Permite realizar nuevas cargas de personal y archivar automáticamente la dotación actual en el histórico.")
        ],
        image_path="screenshots/21_admin_datos.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Administración - Base de Datos",
        transition="fade"
    )
    
    # 26. Administración - Monitoreo de Actividad
    add_content_slide(
        prs,
        title_text="🔐 Administración: Bitácora de Actividad y Auditoría",
        bullets=[
            ("Auditoría en Tiempo Real (Logs)", "Registra cada acción efectuada en la plataforma, indicando el usuario/rol, la acción detallada y la hora exacta."),
            ("Monitoreo de Terminales Activas", "Identifica las computadoras conectadas al servidor en la red local LAN, mostrando su IP y hora de conexión."),
            ("Transparencia Operativa", "Garantiza el control total de los datos al poder auditar quién generó una orden, quién modificó un delito o quién limpió tablas."),
            ("Resiliencia de Conexión", "Monitorea constantemente la conexión con la base de datos SQLite y muestra alertas visuales ante desconexiones.")
        ],
        image_path="screenshots/22_admin_actividad.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Administración - Auditoría de Logs",
        transition="fade"
    )
    
    # 27. Conclusiones y Resumen de Beneficios
    add_content_slide(
        prs,
        title_text="ℹ️ Resumen de Beneficios del Sistema OMAI",
        bullets=[
            ("Agilidad en Toma de Decisiones", "Permite a los comandantes visualizar datos tácticos de forma geográfica e infográfica en tiempo real."),
            ("Estandarización y Calidad", "Garantiza que la correspondencia, partes y órdenes sigan estrictamente las normas doctrinarias de la Armada del Ecuador."),
            ("Ahorro de Tiempo Administrativo", "La importación Excel, la clonación de órdenes y la distribución automatizada liberan al personal de tareas burocráticas."),
            ("Seguridad y Autonomía", "Su diseño offline-first y base de datos local protegen la información en áreas remotas o sin acceso a Internet.")
        ],
        image_path="screenshots/23_acerca_de.png",
        footer_text="SISTEMA OMAI - GT 100.51 | Conclusiones y Autoría",
        transition="zoom" # Animación Zoom para el cierre espectacular
    )
    
    # Guardar Presentación
    output_filename = "SISTEMA_OMAI_Presentacion.pptx"
    prs.save(output_filename)
    print(f"Presentación guardada con éxito como: {output_filename}")
    print(f"Ruta completa: {os.path.abspath(output_filename)}")

if __name__ == '__main__':
    create_presentation()
